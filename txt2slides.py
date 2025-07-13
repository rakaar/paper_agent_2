import sys
import os
import json
import subprocess
import shutil
import tempfile
from pathlib import Path
import time
import argparse
import re

# --- Whitespace Compaction Helper ---
# Keeps single newlines for readability but removes excess internal spaces
# and consecutive blank lines. Also replaces common Unicode space chars.
UNICODE_SPACES = ["\u00A0", "\u2002", "\u2003", "\u2009"]

def compact_whitespace(text: str) -> str:
    """Return `text` with internal whitespace collapsed to save LLM tokens."""
    lines = []
    prev_blank = False
    for raw in text.splitlines():
        # Replace special unicode spaces
        for u in UNICODE_SPACES:
            raw = raw.replace(u, " ")
        stripped = raw.strip()
        # Collapse internal whitespace (incl. tabs) to single space
        stripped = re.sub(r"\s+", " ", stripped)
        if stripped == "":
            if not prev_blank:
                lines.append("")  # keep a single blank line
            prev_blank = True
        else:
            lines.append(stripped)
            prev_blank = False
    return "\n".join(lines).strip()

from pdf2json import call_llm

# Try to import pptx and provide a helpful error message if it's not installed.
try:
    import pptx
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
except ImportError:
    print("Error: The 'python-pptx' library is required. Please install it by running:")
    print("pip install python-pptx")
    sys.exit(1)

# Import SarvamAI for TTS
try:
    from sarvamai import SarvamAI
    from sarvamai.play import save as sarvam_save # Rename to avoid conflict with os.save
except ImportError:
    print("Error: The 'sarvamai' library is required for Text-to-Speech. Please install it by running:")
    print("pip install sarvamai")
    sys.exit(1)

# Sarvam API Key (read from environment variable)
SARVAM_API_KEY = os.getenv("SARVAM_API_KEY")
if not SARVAM_API_KEY:
    print("Error: SARVAM_API_KEY environment variable not set.")
    sys.exit(1)

def create_presentation(slides_data: list, output_filename: str):
    """Creates a PowerPoint presentation from slide data."""
    prs = pptx.Presentation()

    # Find a suitable layout with a body placeholder
    content_layout = None
    for layout_idx, layout in enumerate(prs.slide_layouts):
        for placeholder in layout.placeholders:
            if placeholder.is_placeholder and placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                content_layout = layout
                print(f"Found suitable layout at index {layout_idx} with a BODY placeholder.")
                break
        if content_layout:
            break

    if content_layout is None:
        print("Warning: No slide layout with a BODY placeholder found. Using default layout 0.")
        content_layout = prs.slide_layouts[0] # Fallback to a blank layout

    print(f"Creating presentation '{output_filename}.pptx'...")
    for i, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(content_layout)
        
        # Safely get title and content, providing defaults if keys are missing
        title = slide_info.get("title", f"Slide {i+1}")
        content = slide_info.get("content", "No content provided.")

        # If content is a list, join it into a single string
        if isinstance(content, list):
            content = "\n".join(content)

        # Set the title
        # Assuming the first placeholder is the title for most layouts
        if slide.shapes.title:
            slide.shapes.title.text = title
        else:
            print(f"Warning: Slide {i+1} has no title placeholder. Title not set.")
        
        # Find the body placeholder
        body_shape = None
        for shape in slide.placeholders:
            if shape.is_placeholder and shape.placeholder.format.type == PP_PLACEHOLDER_TYPE.BODY:
                body_shape = shape
                break
        
        if body_shape is None:
            print(f"Warning: Slide {i+1} could not find a suitable body placeholder. Content might not be added.")
            continue

        tf = body_shape.text_frame
        tf.clear() # Clear existing text (like "Click to add text")
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT

    prs.save(f"{output_filename}.pptx")
    print(f"Successfully saved presentation: {output_filename}.pptx")

def save_audio_script(slides_data: list, output_filename: str):
    """Saves the audio narration script to a text file."""
    script_filename = f"{output_filename}_audio_script.txt"
    print(f"Creating audio script '{script_filename}'...")
    with open(script_filename, "w", encoding="utf-8") as f:
        for i, slide_info in enumerate(slides_data):
            slide_num = slide_info.get("slide number", i + 1)
            title = slide_info.get("title", "Untitled Slide")
            audio_text = slide_info.get("audio", "No narration provided.")
            
            f.write(f"--- Slide {slide_num}: {title} ---\n")
            f.write(f"{audio_text}\n\n")
    print(f"Successfully saved audio script: {script_filename}")

def save_slides_json(slides_data: list, output_filename: str):
    """Saves the slides data as a JSON file."""
    json_filename = f"{output_filename}_slides_plan.json"
    print(f"Creating slides plan JSON '{json_filename}'...")
    with open(json_filename, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2)
    print(f"Successfully saved slides plan JSON: {json_filename}")



def fix_json_newlines(json_string: str) -> str:
    """Fixes unescaped newlines inside JSON string values using a robust regex."""
    
    def escape_newlines_in_match(match):
        # The match object gives us the entire string literal, including the quotes.
        # We replace literal newlines with their escaped version inside this matched string.
        return match.group(0).replace('\n', '\\n')

    # This regex robustly finds all JSON string literals. It looks for a double quote,
    # followed by any sequence of characters that are not a backslash or a double quote,
    # or any escaped character (e.g., \", \\, \n), and ends with a double quote.
    string_literal_regex = r'"((?:\\.|[^"\\])*)"'
    
    return re.sub(string_literal_regex, escape_newlines_in_match, json_string, flags=re.DOTALL)

def generate_audio_files(slides_data: list, output_dir: str):
    """Generates audio files from slide narration using Sarvam TTS."""
    print("Generating audio files from slide narration...")
    os.makedirs(output_dir, exist_ok=True)
    client = SarvamAI(api_subscription_key=SARVAM_API_KEY)

    for i, slide_info in enumerate(slides_data):
        audio_text = slide_info.get("audio", "")
        # Use 1-based indexing for slide numbers in filenames
        slide_num = slide_info.get("slide number", i + 1)
        
        if not audio_text:
            print(f"  Skipping audio generation for Slide {slide_num}: No narration text provided.")
            continue

        try:
            print(f"  Generating audio for Slide {slide_num}...")
            audio = client.text_to_speech.convert(
                text=audio_text,
                target_language_code="en-IN",
                model="bulbul:v2",
                speaker="anushka"
            )
            # Save audio files as slide01.wav, slide02.wav, etc.
            audio_file_path = os.path.join(output_dir, f"slide{slide_num:02d}.wav")
            sarvam_save(audio, audio_file_path)
            print(f"  Successfully saved audio for Slide {slide_num} to {audio_file_path}")
        except Exception as e:
            print(f"  Error generating audio for Slide {slide_num}: {e}")

    print("Audio generation complete.")

def create_video_with_ffmpeg(frames_dir: str, audio_dir: str, output_video_path: str):
    """Creates a video from PNG frames and audio files using ffmpeg."""
    print("Creating video using ffmpeg...")
    
    # Check for ffmpeg
    ffmpeg_path = shutil.which("ffmpeg")
    if not ffmpeg_path:
        print("Error: ffmpeg not found. Please install ffmpeg.")
        return

    temp_dir = Path(tempfile.mkdtemp(prefix="video_temp_"))
    try:
        # Add a small delay to ensure files are written to disk
        time.sleep(1)

        if not Path(frames_dir).exists():
            print(f"Error: Frames directory does not exist: {frames_dir}")
            return

        print(f"  Contents of frames directory ({frames_dir}):")
        for item in os.listdir(frames_dir):
            print(f"    - {item}")

        # Glob for PNG files, which marp-cli names `deck.001.png`, `deck.002.png`, etc.
        png_files = sorted(Path(frames_dir).glob("deck.*.png"))
        audio_files = sorted(Path(audio_dir).glob("*.wav"))

        print(f"  PNG files found by glob: {png_files}")

        if not png_files:
            print("Error: No PNG files found in frames directory.")
            return
        if not audio_files:
            print("Error: No audio files found for video creation.")
            return

        # Step 1: Pre-process audio files to a standard format to avoid ffmpeg errors
        print("  Pre-processing audio files...")
        standardized_audio_files = []
        for audio_file in audio_files:
            standardized_path = temp_dir / f"std_{audio_file.name}"
            standardize_cmd = [
                ffmpeg_path,
                "-i", str(audio_file),
                "-acodec", "pcm_s16le", # Standard 16-bit PCM
                "-ar", "44100", # 44.1kHz sample rate
                "-ac", "2", # Stereo
                str(standardized_path)
            ]
            subprocess.run(standardize_cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            standardized_audio_files.append(standardized_path)
        print("  Audio pre-processing complete.")

        # Step 2: Create individual video clips (image + audio) for each slide
        individual_clips = []
        for i, audio_file in enumerate(standardized_audio_files):
            slide_num = i + 1
            if len(png_files) <= i:
                print(f"Warning: Missing PNG for slide {slide_num}. Skipping video creation for this slide.")
                continue
            png_file = png_files[i]
            
            clip_output_path = temp_dir / f"clip_{slide_num:02d}.mp4"
            
            # Get audio duration from the standardized file
            ffprobe_path = shutil.which("ffprobe") or ffmpeg_path # Use ffprobe if available
            duration_cmd = [ffprobe_path, "-i", str(audio_file), "-show_entries", "format=duration", "-v", "quiet", "-of", "csv=p=0"]
            duration_output = subprocess.check_output(duration_cmd).decode("utf-8").strip()
            duration = float(duration_output)

            # Command to create a video clip for one slide
            clip_cmd = [
                ffmpeg_path, 
                "-loop", "1", 
                "-i", str(png_file),
                "-i", str(audio_file),
                "-c:v", "libx264", 
                "-tune", "stillimage", # Optimize for still images
                "-c:a", "aac", 
                "-b:a", "192k", # Audio bitrate
                "-pix_fmt", "yuv420p", 
                "-shortest", 
                "-t", str(duration), # Set video duration to audio duration
                str(clip_output_path)
            ]
            print(f"  Executing ffmpeg clip command: {' '.join(clip_cmd)}")
            subprocess.run(clip_cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            individual_clips.append(clip_output_path)
        
        if not individual_clips:
            print("Error: No individual video clips were created.")
            return

        # Step 2: Concatenate individual video clips
        print("  Concatenating individual video clips...")
        concat_file_list_path = temp_dir / "concat_clips.txt"
        with open(concat_file_list_path, "w", encoding="utf-8") as f:
            for clip_path in individual_clips:
                f.write(f"file '{clip_path}'\n")
        
        final_concat_cmd = [
            ffmpeg_path, 
            "-f", "concat", 
            "-safe", "0", 
            "-i", str(concat_file_list_path),
            "-c", "copy", 
            output_video_path
        ]
        print(f"  Executing ffmpeg concat command: {' '.join(final_concat_cmd)}")
        subprocess.run(final_concat_cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print(f"Successfully created video: {output_video_path}")

    finally:
        shutil.rmtree(temp_dir)

def main():
    """Main function to drive the script."""
    parser = argparse.ArgumentParser(description="Convert text files into a narrated video presentation.")
    parser.add_argument("input_files", nargs='+', help="One or more input text files.")
    parser.add_argument("--figures-path", type=str, help="Optional path to a JSON file containing metadata for figures to be included.")
    parser.add_argument("--max-slides", type=int, help="Desired maximum number of slides (overrides heuristic).")
    parser.add_argument("--slides-only", action="store_true", help="Only generate slides JSON & Markdown; skip audio and video steps.")
    
    args = parser.parse_args()

    input_files = args.input_files
    
    print(f"Reading and combining text from: {', '.join(input_files)}")
    
    combined_text = []
    for file_path in input_files:
        try:
            if file_path.lower().endswith('.pdf'):
                try:
                    import fitz # PyMuPDF
                except ImportError:
                    print("Error: The 'PyMuPDF' library is required to process PDF files. Please install it by running:")
                    print("pip install PyMuPDF")
                    sys.exit(1)
                
                print(f"Extracting text from PDF: {file_path}")
                doc = fitz.open(file_path)
                text = "".join(page.get_text() for page in doc)
                combined_text.append(text)
                doc.close()
            else: # Assume .txt or other text format
                with open(file_path, "r", encoding="utf-8") as f:
                    combined_text.append(f.read())
        except FileNotFoundError:
            print(f"Error: File not found at '{file_path}'. Skipping.")
            continue
    
    if not combined_text:
        print("No valid files were found. Exiting.")
        sys.exit(1)

    full_text = "\n\n---\n\n".join(combined_text)

    # --- Handle Figures --- #
    figures_prompt_injection = ""
    if args.figures_path:
        print(f"Loading figures metadata from: {args.figures_path}")
        try:
            with open(args.figures_path, 'r', encoding='utf-8') as f:
                figures_data = json.load(f)
            
            figures_list_str = ""
            for i, fig in enumerate(figures_data):
                title = fig.get('title', f'Figure {i+1}')
                caption = fig.get('caption', 'No caption available.')
                # The markdown_path will be created by the orchestrator script
                raw_path = fig.get('markdown_path', '')
                # Make the path relative to the location of deck.md (slides/)
                path = os.path.relpath(raw_path, start="slides") if raw_path else ''
                figures_list_str += f"- Figure {i+1}:\n  - Title: {title}\n  - Caption: {caption}\n  - Markdown Path: {path}\n"

            if figures_list_str:
                figures_prompt_injection = f"""\n\n--- AVAILABLE FIGURES ---
You have been provided with a list of figures. Where relevant, you MUST embed these figures into the slide content using their provided Markdown paths (e.g., `![{title}]({path})`).

{figures_list_str}
"""
        except FileNotFoundError:
            print(f"Warning: Figures metadata file not found at '{args.figures_path}'. Continuing without figures.")
        except json.JSONDecodeError:
            print(f"Warning: Could not decode JSON from '{args.figures_path}'. Continuing without figures.")
    # --- End Handle Figures --- #

    # Calculate max_slides based on the number of input files
    if args.max_slides and args.max_slides > 0:
        max_slides = args.max_slides
    else:
        # Heuristic: roughly 1 slide per 1500 chars, capped 15
        max_slides = min(15, max(4, len(full_text)//1500))

    system_prompt = """You are an AI assistant role-playing as a graduate student in a lab meeting, explaining an interesting paper to your peers.
Your tone should be conversational, insightful, and slightly informal. Refer to the paper's authors as 'the authors' or 'the paper,' not 'we'.

Your task is to create a JSON object that represents the slide deck. The JSON object should be a list of slides.
Each slide should have a "slide number", a "title", a "content" field, and an "audio" field.
The "content" field should contain the text for the slide body as a single string, formatted in Markdown.
The "audio" field should contain the narration script for the slide, matching your persona.

IMPORTANT: The output MUST be a single, valid JSON object. Ensure that all strings are properly escaped. For example, use \n for newlines within the content and audio strings, and escape any double quotes. Do not add any extra text or formatting outside of the JSON object itself. The entire response should be parseable by a standard JSON parser."""
    
    # The user_prompt must be defined *after* figures_prompt_injection is created.
    user_prompt = f"""**IMPORTANT:** First, review the list of available figures. You MUST embed these figures in the 'content' of relevant slides.
{figures_prompt_injection}

Now, please break the following text into exactly {max_slides} slides. Each slide must be a JSON object with these exact keys: "slide number", "title", "content", and "audio".
- "slide number": An integer for the slide order.
- "title": A concise title for the slide.
- "content": Keep this extremely minimal:
        * If the slide EMBEDS A FIGURE, use **max 2 short bullet points or <=120 characters**.
        * Otherwise 3-4 bullets or brief paragraph. This is for on-screen text only.
- "audio": This should contain the full, detailed narration for the slide, suitable for text-to-speech. Maximize information transfer here.

Do not include any text, prose, or markdown formatting outside of the main JSON array.

--- TEXT TO CONVERT ---
{full_text}

--- END OF TEXT ---
Remember to include the figures in your response where appropriate."""

    # For transparency, print the full user prompt being sent (can be verbose)
    # Compact whitespace in prompts to save tokens
    compacted_system_prompt = compact_whitespace(system_prompt)
    compacted_user_prompt = compact_whitespace(user_prompt)

    print("\n--- LLM USER PROMPT (truncated to 1500 chars) ---")
    print(user_prompt[:1500] + ("..." if len(user_prompt) > 1500 else ""))
    print("--- END PROMPT ---\n")

    # Save full prompts to a file for user inspection
    debug_prompt_path = Path("slides/full_llm_prompt.txt")
    try:
        debug_prompt_path.parent.mkdir(parents=True, exist_ok=True)
        with open(debug_prompt_path, "w", encoding="utf-8") as f:
            f.write("=== SYSTEM PROMPT (raw) ===\n")
            f.write(system_prompt + "\n\n")
            f.write("=== USER PROMPT (raw) ===\n")
            f.write(user_prompt + "\n\n")
            f.write("=== SYSTEM PROMPT (compacted) ===\n")
            f.write(compacted_system_prompt + "\n\n")
            f.write("=== USER PROMPT (compacted) ===\n")
            f.write(compacted_user_prompt)
        print(f"Full LLM prompt written to {debug_prompt_path}")
    except Exception as e:
        print(f"Warning: could not write debug prompt file: {e}")
    print("Sending content to the language model for processing...")
    try:
        llm_response_str = call_llm(compacted_system_prompt, compacted_user_prompt)
        
        # Parse the JSON response, fixing any unescaped newlines first
        try:
            fixed_llm_response = fix_json_newlines(llm_response_str)
            slides_data = json.loads(fixed_llm_response)
        except json.JSONDecodeError as e:
            print(f"Error decoding JSON even after fixing newlines: {e}")
            # Save the problematic response for inspection
            with open("llm_response_error.json", "w") as f:
                f.write(llm_response_str)
            print("Problematic response saved to llm_response_error.json")
            return
    except Exception as e:
        print(f"An unexpected error occurred during LLM call or processing: {e}")
        sys.exit(1)

    print("LLM response received and parsed successfully.")
    
    # Define output directories and filenames
    # Use only the stem (filename without directories) to avoid nested dirs inside slides/
    base_output_filename = Path(input_files[0]).stem
    slides_dir = Path("slides")
    slides_dir.mkdir(parents=True, exist_ok=True)

    json_output_path = slides_dir / f"{base_output_filename}_slides_plan.json"
    marp_md_path = slides_dir / "deck.md"
    audio_output_dir = slides_dir / "audio"
    frames_output_dir = slides_dir / "frames"
    # Marp CLI needs a dummy file path in the target directory for image sequence output
    frames_output_path_template = frames_output_dir / "deck.png"
    video_output_path = slides_dir / "video.mp4"

    # Ensure output subdir exists (might just be slides/)
    json_output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save slides data as JSON
    with open(json_output_path, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2)
    print(f"Successfully saved slides plan JSON: {json_output_path}")

    # Generate Marp Markdown
    print("Generating Marp Markdown...")
    try:
        subprocess.run(
            [
                sys.executable, 
                "json2marp.py", 
                str(json_output_path), 
                "--out", str(marp_md_path)
            ],
            check=True
        )
        print(f"Successfully created Marp Markdown: {marp_md_path}")
    except subprocess.CalledProcessError as e:
        print(f"Error generating Marp Markdown: {e}")
        print(f"Stdout: {e.stdout.decode()}")
        print(f"Stderr: {e.stderr.decode()}")
        sys.exit(1)

    # Generate audio files (skip in slides-only mode)
    if not args.slides_only:
        generate_audio_files(slides_data, str(audio_output_dir)) 

    # Render Marp Markdown to PNG frames
    print("Rendering Marp Markdown to PNG frames...")
    frames_output_dir.mkdir(parents=True, exist_ok=True) # Ensure frames directory exists
    try:
        # Use npx to run marp-cli
        subprocess.run(
            [
                "npx", "marp", 
                str(marp_md_path),
                "--images", "png",
                "--image-scale", "2",
                "--allow-local-files",
                # Provide a dummy file path; marp-cli will use its basename for the sequence
                "--output", str(frames_output_path_template)
            ],
            check=True
        )
        print(f"Successfully rendered PNG frames to: {frames_output_dir}")
    except subprocess.CalledProcessError as e:
        print(f"Error rendering Marp Markdown to PNGs: {e}")
        print(f"Stdout: {e.stdout.decode()}")
        print(f"Stderr: {e.stderr.decode()}")
        sys.exit(1)

    if args.slides_only:
        print("Slides-only mode: frames rendered. Skipping audio/video generation.")
        return

    # Build video from PNG frames and audio
    print("Building video from PNG frames and audio...")
    try:
        create_video_with_ffmpeg(str(frames_output_dir), str(audio_output_dir), str(video_output_path))
        print(f"Successfully created video: {video_output_path}")
    except Exception as e:
        print(f"Error building video: {e}")
        if isinstance(e, subprocess.CalledProcessError):
            print(f"ffmpeg stdout: {e.stdout.decode() if e.stdout else 'N/A'}")
            print(f"ffmpeg stderr: {e.stderr.decode() if e.stderr else 'N/A'}")
        sys.exit(1)

    print("\nAll tasks completed.")

if __name__ == "__main__":
    main()